from datetime import datetime, timezone, timedelta

import PyPDF2
import torch
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_qdrant import Qdrant
import sys
from langchain_text_splitters import TokenTextSplitter
from opensearchpy import OpenSearch
from pptx import Presentation
from qdrant_client import QdrantClient

from qdrant_client.http import models
import docx
import os
from dotenv import load_dotenv

def get_files(dir):
    file_list = []
    for dir, _, filenames in os.walk(dir):
        for f in filenames:
            file_list.append(os.path.join(dir, f))
    return file_list


def getTextAndMetadataFromWord(filename):
    doc = docx.Document(filename)

    # Extract text
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    text_content = '\n'.join(fullText)

    # Extract metadata
    properties = doc.core_properties

    # Get file system metadata
    file_stats = os.stat(filename)

    metadata = {
        'Author': properties.author or 'Unknown',
        'Title': properties.title or 'Untitled',
        'Subject': properties.subject or 'No Subject',
        'Keywords': properties.keywords or 'No Keywords',
        'Created': properties.created.isoformat() if properties.created else 'Unknown',
        'Modified': properties.modified.isoformat() if properties.modified else 'Unknown',
        'Last_Modified_By': properties.last_modified_by or 'Unknown',
        'Revision': properties.revision or 'Unknown',
        'File_Size': f"{file_stats.st_size / (1024 * 1024):.2f} MB",
        'Creation_Date': datetime.fromtimestamp(file_stats.st_ctime).isoformat(),
        'File_Modified': datetime.fromtimestamp(file_stats.st_mtime).isoformat(),
        'File_Accessed': datetime.fromtimestamp(file_stats.st_atime).isoformat(),
    }

    return text_content, metadata


def getTextAndMetadataFromPPTX(filename):
    prs = Presentation(filename)

    # Extract text
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                fullText.append(shape.text)
    text_content = '\n'.join(fullText)

    # Extract metadata
    core_properties = prs.core_properties

    # Get file system metadata
    file_stats = os.stat(filename)

    metadata = {
        'Author': core_properties.author or 'Unknown',
        'Title': core_properties.title or 'Untitled',
        'Subject': core_properties.subject or 'No Subject',
        'Keywords': core_properties.keywords or 'No Keywords',
        'Created': core_properties.created.isoformat() if core_properties.created else 'Unknown',
        'Modified': core_properties.modified.isoformat() if core_properties.modified else 'Unknown',
        'Last_Modified_By': core_properties.last_modified_by or 'Unknown',
        'Revision': core_properties.revision or 'Unknown',
        'Category': core_properties.category or 'Uncategorized',
        'Comments': core_properties.comments or 'No Comments',
        'Content_Status': core_properties.content_status or 'Unknown',
        'Slide_Count': len(prs.slides),
        'File_Size': f"{file_stats.st_size / (1024 * 1024):.2f} MB",
        'Creation_Date': datetime.fromtimestamp(file_stats.st_ctime).isoformat(),
        'File_Modified': datetime.fromtimestamp(file_stats.st_mtime).isoformat(),
        'File_Accessed': datetime.fromtimestamp(file_stats.st_atime).isoformat(),
    }

    return text_content, metadata


def parse_pdf_date(date_string):
    if date_string == 'Unknown' or not date_string:
        print(f"Date string is empty or unknown: {date_string}")
        return datetime.now().isoformat()

    # Remove the 'D:' prefix if it exists
    if date_string.startswith('D:'):
        date_string = date_string[2:]

    # Parse the date string
    try:
        # Parse the main part of the date
        dt = datetime.strptime(date_string[:14], '%Y%m%d%H%M%S')
    except ValueError as ve:
        print(f"Error processing {date_string}: {str(ve)}")
        return datetime.now().isoformat()

    try:
        # Handle the timezone information
        if len(date_string) > 14:
            tz_string = date_string[14:]
            if tz_string[0] in ['+', '-']:
                tz_hours = int(tz_string[1:3])
                tz_minutes = int(tz_string[4:6])
                tz_offset = (tz_hours * 60 + tz_minutes) * 60
                if tz_string[0] == '-':
                    tz_offset = -tz_offset
                dt = dt.replace(tzinfo=timezone(timedelta(seconds=tz_offset)))
            elif tz_string == 'Z':
                dt = dt.replace(tzinfo=timezone.utc)

        return dt.isoformat()
    except Exception as exc:
        print(f"Error processing {date_string}: {str(exc)}")
        pass

    return dt.isoformat()

def getTextAndMetadataFromPDF(filename):
    file_content = ""
    metadata = {}

    try:
        reader = PyPDF2.PdfReader(filename)

        # Extract text content
        for page in reader.pages:
            file_content += " " + page.extract_text()

        # Extract metadata
        info = reader.metadata
        if info:
            metadata = {
                'Author': info.get('/Author', 'Unknown'),
                'Creator': info.get('/Creator', 'Unknown'),
                'Producer': info.get('/Producer', 'Unknown'),
                'Subject': info.get('/Subject', 'No Subject'),
                'Title': info.get('/Title', 'Untitled'),
                'Creation_Date': parse_pdf_date(info.get('/CreationDate', 'Unknown')),
                'Modification_Date': parse_pdf_date(info.get('/ModDate', 'Unknown')),
            }

        # Add page count
        metadata['Page_Count'] = len(reader.pages)

        # Get file system metadata
        file_stats = os.stat(filename)
        metadata.update({
            'File_Size': f"{file_stats.st_size / (1024 * 1024):.2f} MB",
            'File_Created': datetime.fromtimestamp(file_stats.st_ctime).strftime('%Y-%m-%d %H:%M:%S'),
            'File_Modified': datetime.fromtimestamp(file_stats.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
            'File_Accessed': datetime.fromtimestamp(file_stats.st_atime).strftime('%Y-%m-%d %H:%M:%S'),
        })

    except Exception as exc:
        print(f"Error processing {filename}: {str(exc)}")
        file_content = "Empty due to extraction error."
        metadata['Error'] = str(exc)

    return file_content.strip(), metadata

def main_indexing(mypath):
    load_dotenv()
    embedding_model = os.getenv("EMBEDDING_MODEL")
    open_search_ip = os.getenv("OPENSEARCH_IP")
    open_search_user = os.getenv("OPENSEARCH_USER")
    open_search_pass = os.getenv("OPENSEARCH_PASSWORD")
    open_search_port = os.getenv("OPENSEARCH_PORT")
    qdrant_ip = os.getenv("QDRANT_IP")
    qdrant_port = os.getenv("QDRANT_PORT")
    qdrant_key = os.getenv("QDRANT_API")
    INDEX_NAME_SEMANTIC = os.getenv("INDEX_NAME_SEMANTIC")
    INDEX_NAME_LEXICAL = os.getenv("INDEX_NAME_LEXICAL")

    model_name = embedding_model
    if torch.cuda.is_available():
        model_kwargs = {'device': 'cuda'}
    elif torch.backends.mps.is_available():
        model_kwargs = {'device': 'mps'}
    else:
        model_kwargs = {'device': 'cpu'}
    encode_kwargs = {'normalize_embeddings': True}
    hf = HuggingFaceEmbeddings(
        model_name=model_name,
        model_kwargs=model_kwargs,
        encode_kwargs=encode_kwargs
    )
    qdrant_client = QdrantClient(qdrant_ip, port=qdrant_port, timeout=500, api_key=qdrant_key,prefer_grpc=True,verify=False,https=False)
    collections = qdrant_client.get_collections()
    collection_exists = any(collection.name == INDEX_NAME_SEMANTIC for collection in collections.collections)

    # If the collection exists, delete it
    if collection_exists:
        qdrant_client.delete_collection(collection_name=INDEX_NAME_SEMANTIC)
        print(f"Collection '{INDEX_NAME_SEMANTIC}' deleted.")

    response = qdrant_client.create_collection(
        collection_name=INDEX_NAME_SEMANTIC,
        vectors_config=models.VectorParams(size=768, distance=models.Distance.DOT),
        hnsw_config=models.HnswConfigDiff(max_indexing_threads=1),
        optimizers_config=models.OptimizersConfigDiff(memmap_threshold=20000, max_optimization_threads=4),
        quantization_config=models.ScalarQuantization(
            scalar=models.ScalarQuantizationConfig(
                type=models.ScalarType.INT8,
                always_ram=True,
            ),
        ),
        shard_number=4
    )


    qdrant = Qdrant(qdrant_client, INDEX_NAME_SEMANTIC, hf)

    print("Qdrant Response = ", response)
    host = open_search_ip
    port = open_search_port
    auth = (open_search_user, open_search_pass)  # ('admin', 'admin')

    # Create the client with SSL/TLS and hostname verification disabled.

    open_search_client = OpenSearch(
        hosts=[{'host': host, 'port': port}],
        http_auth=auth,
        use_ssl=False,
        verify_certs=False,
        ssl_assert_hostname=False,
        ssl_show_warn=False,
        timeout=500,
        max_retries=10
        # connection_class=RequestsHttpConnection
        #    http_compress = True, # enables gzip compression for request bodies
        #    use_ssl = False,
        #   verify_certs = False,
        #    ssl_assert_hostname = False,
        #    ssl_show_warn = False
    )
    print("Open Search Client created...")
    index_body = {
        "settings": {
            "index": {
                "refresh_interval": -1,
                "number_of_shards": 4,
                "number_of_replicas": 0
            }
        },
        "mappings": {
            "properties": {
                "pmid": {"type": "keyword"},
                "location": {"type": "text"},
                "full_text": {"type": "text"},
                "journal": {"type": "text", "index": False},
                "pubdate": {"type": "date"},
                "authors": {
                    "type": "nested",
                    "properties": {
                        "lastname": {"type": "text"},
                        "forename": {"type": "text"},
                        "initials": {"type": "text"},
                        "identifier": {"type": "text"},
                        "affiliation": {"type": "text"},
                    }
                }
            }
        }
    }

    index_exists = open_search_client.indices.exists(index=INDEX_NAME_LEXICAL)
    if index_exists:
        # If the index exists, delete it
        open_search_client.indices.delete(index=INDEX_NAME_LEXICAL)
        print(f"Index '{INDEX_NAME_LEXICAL}' deleted.")

    response = open_search_client.indices.create(INDEX_NAME_LEXICAL, body=index_body)
    print(response)

    print("Indexing...")
    onlyfiles = get_files(mypath)
    file_content = ""
    for file in onlyfiles:
        file_content = ""
        meta = {}
        if file.find("~") > 0:
            file_content = "Empty due to ~ in file name."
            print("Document title with ~: " + file)
        elif file.endswith(".pdf"):
            try:
                print("indexing " + file)
                file_content,meta = getTextAndMetadataFromPDF(file)
            except Exception as exc:
                file_content = "Empty due to extraction error."
        elif file.endswith(".txt") or file.endswith(".md") or file.endswith(".markdown"):
            print("indexing " + file)
            f = open(file, 'r', encoding='utf-8', errors='ignore')
            file_content = f.read()
            file_stats = os.stat(file)
            meta = {
                'File_Size': f"{file_stats.st_size / (1024 * 1024):.2f} MB",
                'Creation_Date': datetime.fromtimestamp(file_stats.st_ctime).isoformat(),
                'File_Modified': datetime.fromtimestamp(file_stats.st_mtime).isoformat(),
                'File_Accessed': datetime.fromtimestamp(file_stats.st_atime).isoformat(),
            }

            f.close()
        elif file.endswith(".docx"):
            print("indexing " + file)
            file_content,meta = getTextAndMetadataFromWord(file)
        elif file.endswith(".pptx"):
            print("indexing " + file)
            file_content,meta = getTextAndMetadataFromPPTX(file)
        else:
            continue
        text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=50)
        texts = text_splitter.split_text(file_content)
        metadata = []
        for i in range(0, len(texts)):
            metadata.append({"location": file})
        qdrant.add_texts(texts, metadatas=metadata)
        doc = {
            "pmid": "",
            "location": file,
            "full_text": file_content,
            "journal": "",
            "pubdate": str(meta['Creation_Date']),
        }
        try:
            open_search_client.index(INDEX_NAME_LEXICAL,body=doc)
            print(f"Document indexed successfully "+file)
        except Exception as e:
            print(f"Error indexing document: {str(e)}")
    open_search_client.indices.refresh(index=INDEX_NAME_LEXICAL)
    count = open_search_client.count(index=INDEX_NAME_LEXICAL)
    print(f"Count of documents indexed in Open Search: {count['count']}")



if __name__ == "__main__":
    arguments = sys.argv
    if len(arguments)>1:
        main_indexing(arguments[1])
    else:
        print("You need to provide a path to folder with documents to index as command line argument")